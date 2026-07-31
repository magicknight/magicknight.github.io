#!/usr/bin/env python3
"""Build the public ScanLine Gemini Startup Forum pitch deck."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
ASSETS = ROOT / "assets"
OUT = ROOT / "ScanLine_AI_Gemini_Startup_Forum_2026.pptx"

W = 13.333
H = 7.5

NAVY = RGBColor(7, 26, 43)
NAVY_2 = RGBColor(13, 45, 70)
INK = RGBColor(11, 23, 39)
INK_SOFT = RGBColor(38, 56, 76)
BLUE = RGBColor(47, 109, 246)
TEAL = RGBColor(24, 184, 170)
TEAL_DARK = RGBColor(8, 126, 118)
MINT = RGBColor(103, 228, 214)
AMBER = RGBColor(241, 166, 58)
PAPER = RGBColor(245, 248, 251)
WHITE = RGBColor(255, 255, 255)
MUTED = RGBColor(99, 116, 135)
LINE = RGBColor(219, 229, 238)
PALE_BLUE = RGBColor(234, 240, 255)
PALE_TEAL = RGBColor(223, 246, 243)
PALE_AMBER = RGBColor(255, 241, 216)

FONT = "Aptos"
FONT_DISPLAY = "Aptos Display"


def add_rect(
    slide,
    x,
    y,
    w,
    h,
    fill,
    *,
    radius=False,
    line=None,
    line_width=1,
    transparency=0,
):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(
        shape_type, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.fill.transparency = transparency
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_width)
    if radius:
        shape.adjustments[0] = 0.13
    return shape


def add_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    *,
    size=18,
    color=INK,
    bold=False,
    font=FONT,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin=0,
    fit=False,
    line_spacing=1.0,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    if fit:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    return box


def add_rich_text(
    slide,
    runs,
    x,
    y,
    w,
    h,
    *,
    size=18,
    color=INK,
    valign=MSO_ANCHOR.TOP,
    align=PP_ALIGN.LEFT,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    for item in runs:
        r = p.add_run()
        r.text = item["text"]
        r.font.name = item.get("font", FONT)
        r.font.size = Pt(item.get("size", size))
        r.font.bold = item.get("bold", False)
        r.font.color.rgb = item.get("color", color)
    return box


def add_image_crop(slide, path, x, y, w, h):
    from PIL import Image

    path = Path(path)
    with Image.open(path) as img:
        iw, ih = img.size
    target_ratio = w / h
    image_ratio = iw / ih
    pic = slide.shapes.add_picture(
        str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h)
    )
    if image_ratio > target_ratio:
        visible_fraction = target_ratio / image_ratio
        crop = (1 - visible_fraction) / 2
        pic.crop_left = crop
        pic.crop_right = crop
    else:
        visible_fraction = image_ratio / target_ratio
        crop = (1 - visible_fraction) / 2
        pic.crop_top = crop
        pic.crop_bottom = crop
    return pic


def add_circle(slide, x, y, d, fill, line=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
    else:
        shape.line.fill.background()
    return shape


def add_line(slide, x1, y1, x2, y2, color=LINE, width=1.5):
    line = slide.shapes.add_connector(
        1, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    line.line.color.rgb = color
    line.line.width = Pt(width)
    return line


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_brand(slide, dark=False):
    fg = WHITE if dark else NAVY
    add_circle(slide, 0.42, 0.3, 0.27, fg)
    add_circle(slide, 0.485, 0.365, 0.14, NAVY if dark else WHITE)
    add_text(
        slide,
        "SCANLINE",
        0.79,
        0.27,
        1.35,
        0.28,
        size=10,
        color=fg,
        bold=True,
    )


def add_footer(slide, number, *, dark=False, note=None):
    color = RGBColor(143, 167, 184) if dark else MUTED
    if note:
        add_text(slide, note, 0.55, 7.15, 10.9, 0.18, size=7.5, color=color)
    add_text(
        slide,
        f"{number:02d}",
        12.25,
        7.08,
        0.48,
        0.23,
        size=8,
        color=color,
        bold=True,
        align=PP_ALIGN.RIGHT,
    )


def add_kicker(slide, text, *, dark=False):
    color = MINT if dark else TEAL_DARK
    add_rect(slide, 0.62, 0.85, 0.32, 0.025, color)
    add_text(slide, text.upper(), 1.02, 0.69, 4.4, 0.3, size=9, color=color, bold=True)


def add_title(slide, title, subtitle=None, *, dark=False, size=31):
    color = WHITE if dark else INK
    add_text(
        slide,
        title,
        0.62,
        1.08,
        12.0,
        0.9,
        size=size,
        color=color,
        bold=True,
        font=FONT_DISPLAY,
        fit=True,
    )
    if subtitle:
        sub_color = RGBColor(178, 197, 210) if dark else MUTED
        add_text(slide, subtitle, 0.65, 1.93, 11.3, 0.52, size=14, color=sub_color)


def add_pill(slide, text, x, y, w, *, fill=PALE_TEAL, color=TEAL_DARK):
    add_rect(slide, x, y, w, 0.34, fill, radius=True)
    add_text(
        slide,
        text,
        x,
        y + 0.03,
        w,
        0.23,
        size=8.5,
        color=color,
        bold=True,
        align=PP_ALIGN.CENTER,
    )


def add_card(slide, x, y, w, h, *, fill=WHITE, line=LINE):
    return add_rect(slide, x, y, w, h, fill, radius=True, line=line, line_width=0.8)


def add_bullet_list(slide, items, x, y, w, h, *, color=INK_SOFT, size=14):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(9)
        p.level = 0
        p.text = f"•  {item}"
    return box


def slide_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)
    add_brand(slide, dark=True)
    add_rect(slide, 8.28, 0, 5.053, 7.5, NAVY_2)
    add_image_crop(slide, ASSETS / "bg.jpg", 8.45, 0.44, 4.42, 6.62)
    add_pill(
        slide,
        "GEMINI STARTUP FORUM 2026",
        0.68,
        1.36,
        2.35,
        fill=RGBColor(16, 67, 86),
        color=MINT,
    )
    add_text(
        slide,
        "The intelligence layer for\nproduction-line X-ray CT",
        0.65,
        2.02,
        7.4,
        2.0,
        size=38,
        color=WHITE,
        bold=True,
        font=FONT_DISPLAY,
        line_spacing=0.92,
    )
    add_text(
        slide,
        "Physics-informed AI for fast, trustworthy 3D quality decisions.",
        0.69,
        4.28,
        6.6,
        0.58,
        size=17,
        color=RGBColor(190, 211, 224),
    )
    add_line(slide, 0.68, 5.25, 7.35, 5.25, RGBColor(48, 81, 103), 1)
    add_text(
        slide,
        "Zhihua Liang  ·  Co-founder",
        0.69,
        5.55,
        3.8,
        0.3,
        size=10,
        color=MINT,
        bold=True,
    )
    add_text(
        slide,
        "scanline.cn",
        0.69,
        5.93,
        2.0,
        0.3,
        size=10,
        color=WHITE,
        bold=True,
    )
    add_text(
        slide,
        "Concept visualization",
        10.6,
        6.74,
        1.9,
        0.18,
        size=7,
        color=RGBColor(210, 225, 234),
        align=PP_ALIGN.RIGHT,
    )
    return slide


def slide_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, PAPER)
    add_brand(slide)
    add_kicker(slide, "The structural problem")
    add_title(
        slide,
        "Manufacturing inspection still trades depth for speed.",
        "Internal defects demand 3D evidence, but production decisions happen at takt time.",
    )
    cards = [
        (
            "MANUAL / DESTRUCTIVE",
            "Slow feedback",
            "Operator variability and sampled evidence leave most production unseen.",
            AMBER,
        ),
        (
            "OFFLINE INDUSTRIAL CT",
            "Rich, but late",
            "Full 3D evidence often arrives after the process can still be corrected.",
            BLUE,
        ),
        (
            "FAST 2D X-RAY",
            "Fast, but ambiguous",
            "Occlusion and depth mixing can hide or mislocalize internal defects.",
            TEAL,
        ),
    ]
    for i, (label, headline, body, accent) in enumerate(cards):
        x = 0.65 + i * 4.17
        add_card(slide, x, 2.68, 3.78, 2.52)
        add_rect(slide, x, 2.68, 3.78, 0.09, accent)
        add_text(slide, label, x + 0.25, 2.98, 3.2, 0.24, size=8.5, color=accent, bold=True)
        add_text(
            slide,
            headline,
            x + 0.25,
            3.38,
            3.2,
            0.78,
            size=23,
            color=INK,
            bold=True,
            font=FONT_DISPLAY,
        )
        add_text(slide, body, x + 0.25, 4.22, 3.15, 0.72, size=12.3, color=MUTED)
    add_rect(slide, 0.65, 5.56, 12.0, 0.91, PALE_TEAL, radius=True)
    add_rich_text(
        slide,
        [
            {"text": "OUR THESIS  ", "bold": True, "color": TEAL_DARK, "size": 10},
            {
                "text": "Co-design acquisition, reconstruction, and the quality decision—then optimize the whole loop for the defect that must be caught.",
                "bold": True,
                "color": INK,
                "size": 16,
            },
        ],
        0.98,
        5.85,
        11.2,
        0.35,
        valign=MSO_ANCHOR.MIDDLE,
    )
    add_footer(slide, 2)
    return slide


def slide_solution(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)
    add_brand(slide, dark=True)
    add_kicker(slide, "The ScanLine loop", dark=True)
    add_title(
        slide,
        "From X-ray projections to an auditable action.",
        "AI connects sparse measurements, physical consistency, defect evidence, uncertainty, and factory feedback.",
        dark=True,
    )
    steps = [
        ("01", "Part + defect truth", "Freeze the decision and its error costs"),
        ("02", "Sparse acquisition", "Measure what the decision needs"),
        ("03", "Physics-informed AI", "Infer under projection constraints"),
        ("04", "Defect decision", "Localize, score, and expose confidence"),
        ("05", "Factory feedback", "Rescan, review, and update process"),
    ]
    for i, (num, name, desc) in enumerate(steps):
        x = 0.62 + i * 2.51
        add_card(
            slide,
            x,
            2.73,
            2.22,
            2.44,
            fill=RGBColor(14, 48, 72),
            line=RGBColor(44, 82, 105),
        )
        add_text(slide, num, x + 0.2, 2.99, 0.5, 0.25, size=9, color=MINT, bold=True)
        add_text(
            slide,
            name,
            x + 0.2,
            3.42,
            1.82,
            0.65,
            size=17,
            color=WHITE,
            bold=True,
            font=FONT_DISPLAY,
        )
        add_text(slide, desc, x + 0.2, 4.27, 1.8, 0.6, size=10.7, color=RGBColor(167, 190, 205))
        if i < 4:
            add_text(slide, "→", x + 2.24, 3.63, 0.27, 0.3, size=15, color=MINT, bold=True)
    metrics = [
        "cycle time",
        "minimum defect",
        "recall",
        "false-positive rate",
        "dose",
        "calibration",
    ]
    for i, metric in enumerate(metrics):
        add_pill(
            slide,
            metric.upper(),
            0.68 + i * 2.03,
            5.75,
            1.78,
            fill=RGBColor(16, 67, 86),
            color=RGBColor(181, 233, 226),
        )
    add_text(
        slide,
        "Joint acceptance—not image quality alone",
        4.75,
        6.42,
        3.85,
        0.28,
        size=10,
        color=RGBColor(168, 193, 207),
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, 3, dark=True)
    return slide


def slide_ai(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    add_brand(slide)
    add_kicker(slide, "Why AI is core")
    add_title(
        slide,
        "The product is a decision system—not a prettier reconstruction.",
        "Three synchronized AI layers compress time while keeping the result testable.",
    )
    layers = [
        (
            "1",
            "Acquisition intelligence",
            "Choose views and geometry around the defect task.",
            BLUE,
            PALE_BLUE,
        ),
        (
            "2",
            "Physics-informed inference",
            "Use learned priors without abandoning projection evidence.",
            TEAL_DARK,
            PALE_TEAL,
        ),
        (
            "3",
            "Decision + uncertainty",
            "Return defect evidence, confidence, and a safe fallback.",
            RGBColor(173, 101, 12),
            PALE_AMBER,
        ),
    ]
    for i, (num, title, body, accent, pale) in enumerate(layers):
        x = 0.72 + i * 4.18
        add_card(slide, x, 2.65, 3.78, 2.27, fill=pale, line=pale)
        add_circle(slide, x + 0.28, 2.95, 0.45, accent)
        add_text(
            slide,
            num,
            x + 0.28,
            3.06,
            0.45,
            0.18,
            size=10,
            color=WHITE,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_text(slide, title, x + 0.28, 3.55, 3.08, 0.72, size=18, color=INK, bold=True)
        add_text(slide, body, x + 0.28, 4.36, 3.05, 0.55, size=12, color=INK_SOFT)
    add_line(slide, 2.64, 5.55, 10.69, 5.55, LINE, 2)
    loop = [
        ("MEASURE", BLUE),
        ("INFER", TEAL),
        ("DECIDE", AMBER),
        ("LEARN", BLUE),
    ]
    for i, (text_value, color) in enumerate(loop):
        x = 1.35 + i * 3.04
        add_circle(slide, x, 5.27, 0.58, color)
        add_text(
            slide,
            text_value,
            x - 0.25,
            6.0,
            1.08,
            0.24,
            size=9,
            color=color,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
    add_text(
        slide,
        "Low confidence triggers rescan or human review.",
        4.17,
        6.48,
        5.0,
        0.28,
        size=11,
        color=MUTED,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, 4)
    return slide


def slide_proof(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, PAPER)
    add_brand(slide)
    add_kicker(slide, "Co-founder proof, clean product boundary")
    add_title(
        slide,
        "We know what it takes to move CT from the lab to the line.",
        "Prior operating experience is evidence. ScanLine’s product and IP path are new work.",
    )
    add_card(slide, 0.67, 2.63, 6.25, 3.95, fill=NAVY, line=NAVY)
    add_image_crop(slide, ASSETS / "deltaray-line-demo.png", 0.82, 2.8, 5.95, 3.42)
    add_rect(slide, 0.82, 5.33, 5.95, 0.89, NAVY, transparency=15)
    add_text(
        slide,
        "Prior DeltaRAY inline deployment",
        1.04,
        5.51,
        3.5,
        0.24,
        size=11,
        color=WHITE,
        bold=True,
    )
    add_text(
        slide,
        "Zhihua’s prior work · not a ScanLine product",
        1.04,
        5.87,
        4.2,
        0.2,
        size=8.5,
        color=MINT,
    )
    add_pill(slide, "ESTABLISHED", 7.4, 2.72, 1.38)
    add_text(
        slide,
        "Zhihua previously led AI back-end, reconstruction, and system integration for the DeltaRAY X100 inline CT platform.",
        7.4,
        3.2,
        5.08,
        1.02,
        size=18,
        color=INK,
        bold=True,
        font=FONT_DISPLAY,
    )
    add_line(slide, 7.4, 4.47, 12.43, 4.47, LINE, 1)
    add_pill(slide, "SCANLINE · IN PROGRESS", 7.4, 4.8, 2.05, fill=PALE_BLUE, color=BLUE)
    add_text(
        slide,
        "Independent acquisition design, code, datasets, product architecture, and IP route.",
        7.4,
        5.32,
        4.95,
        0.66,
        size=15,
        color=INK_SOFT,
        bold=True,
    )
    add_text(
        slide,
        "The pictured product and associated IP belong to DeltaRAY.",
        7.4,
        6.21,
        4.85,
        0.3,
        size=9,
        color=MUTED,
    )
    add_footer(
        slide,
        5,
        note="Boundary note: Zhihua’s prior DeltaRAY work demonstrates co-founder capability; no DeltaRAY product, customer, or IP is claimed by ScanLine.",
    )
    return slide


def slide_state(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    add_brand(slide)
    add_kicker(slide, "Current state")
    add_title(
        slide,
        "The ambition stays bold because the open bridge is explicit.",
        "We separate what is established, what is being built, and what must be measured next.",
    )
    states = [
        (
            "ESTABLISHED",
            PALE_TEAL,
            TEAL_DARK,
            "Inline CT system experience",
            [
                "AI back-end and reconstruction",
                "System integration and line-side deployment",
                "Understanding of industrial acceptance",
            ],
        ),
        (
            "IN PROGRESS",
            PALE_BLUE,
            BLUE,
            "Independent ScanLine stack",
            [
                "Acquisition and architecture design",
                "Physics-informed reconstruction + decision",
                "Data, validation, and clean IP contracts",
            ],
        ),
        (
            "OPEN BRIDGE",
            PALE_AMBER,
            RGBColor(169, 96, 0),
            "Production-speed acceptance",
            [
                "2–5 s target for a defined workpiece",
                "Defect sensitivity at required recall / FPR",
                "Dose, drift, and uncertainty performance",
            ],
        ),
    ]
    for i, (tag, pale, accent, headline, bullets) in enumerate(states):
        x = 0.65 + i * 4.18
        add_card(slide, x, 2.68, 3.78, 3.42, fill=WHITE, line=LINE)
        add_rect(slide, x, 2.68, 3.78, 0.11, accent)
        add_pill(slide, tag, x + 0.25, 3.02, 1.4 if i != 1 else 1.5, fill=pale, color=accent)
        add_text(
            slide,
            headline,
            x + 0.25,
            3.62,
            3.2,
            0.58,
            size=21,
            color=INK,
            bold=True,
            font=FONT_DISPLAY,
        )
        add_bullet_list(slide, bullets, x + 0.25, 4.47, 3.05, 1.35, size=11.2)
    add_rect(slide, 2.15, 6.39, 9.03, 0.46, NAVY, radius=True)
    add_text(
        slide,
        "TARGET  ·  Full 3D quality decisions at production speed for a qualified inspection task",
        2.15,
        6.51,
        9.03,
        0.21,
        size=10,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, 6)
    return slide


def slide_probe(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)
    add_brand(slide, dark=True)
    add_kicker(slide, "The decisive 90-day probe", dark=True)
    add_title(
        slide,
        "One part. One defect. One line.",
        "A narrow wedge makes the hardest claim falsifiable—and creates a repeatable deployment template if it wins.",
        dark=True,
    )
    add_rect(slide, 0.68, 2.66, 12.0, 0.88, RGBColor(16, 55, 79), radius=True)
    for i, (big, small) in enumerate(
        [
            ("1 PART", "stable geometry"),
            ("1 DEFECT", "qualified truth"),
            ("1 LINE", "real takt time"),
        ]
    ):
        x = 1.0 + i * 4.05
        add_text(slide, big, x, 2.87, 1.48, 0.26, size=13, color=MINT, bold=True)
        add_text(slide, small, x + 1.52, 2.87, 1.75, 0.26, size=11, color=WHITE)
    phases = [
        (
            "DAYS 0–30",
            "FREEZE",
            "Part, defect taxonomy, truth method, dense baseline, and joint acceptance metrics.",
        ),
        (
            "DAYS 31–60",
            "BUILD",
            "Sparse/dense pairs, physics-informed pipeline, and calibrated uncertainty.",
        ),
        (
            "DAYS 61–90",
            "DECIDE",
            "Blind comparison: deploy, repair the observable, or stop that wedge.",
        ),
    ]
    for i, (days, verb, body) in enumerate(phases):
        x = 0.68 + i * 4.18
        add_card(
            slide,
            x,
            4.0,
            3.78,
            2.03,
            fill=RGBColor(14, 48, 72),
            line=RGBColor(44, 82, 105),
        )
        add_text(slide, days, x + 0.24, 4.26, 1.8, 0.23, size=9, color=MINT, bold=True)
        add_text(slide, verb, x + 0.24, 4.69, 1.7, 0.38, size=22, color=WHITE, bold=True)
        add_text(slide, body, x + 0.24, 5.2, 3.16, 0.57, size=10.8, color=RGBColor(170, 192, 206))
    add_text(
        slide,
        "SUCCESS = faster cycle time at the same qualified defect threshold—not a nicer-looking image.",
        1.2,
        6.5,
        10.9,
        0.28,
        size=11,
        color=MINT,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, 7, dark=True)
    return slide


def slide_wedge(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, PAPER)
    add_brand(slide)
    add_kicker(slide, "Commercial entry")
    add_title(
        slide,
        "Start where hidden defects are costly and truth is obtainable.",
        "Candidate sectors—not customer or deployment claims. The first wedge will be selected by evidence access and economics.",
    )
    cards = [
        (
            "battery.jpg",
            "Battery cells",
            "Internal alignment, folds, contamination, and structural anomalies.",
            "High value · demanding cycle time",
        ),
        (
            "auto.jpg",
            "Castings + welds",
            "Porosity, cracks, inclusions, and internal joining defects.",
            "Clear defect economics · strong CT fit",
        ),
        (
            "chips.jpg",
            "Electronics",
            "Hidden interconnect, packaging, and assembly anomalies.",
            "Small features · high density",
        ),
    ]
    for i, (img, title, body, fit) in enumerate(cards):
        x = 0.65 + i * 4.18
        add_card(slide, x, 2.63, 3.78, 3.55, fill=WHITE, line=LINE)
        add_image_crop(slide, ASSETS / img, x + 0.06, 2.69, 3.66, 1.52)
        add_text(slide, title, x + 0.25, 4.45, 3.15, 0.35, size=20, color=INK, bold=True)
        add_text(slide, body, x + 0.25, 4.94, 3.12, 0.55, size=11.2, color=MUTED)
        add_pill(slide, fit.upper(), x + 0.25, 5.64, 2.92, fill=PALE_TEAL, color=TEAL_DARK)
    add_text(
        slide,
        "Selection gate",
        0.75,
        6.56,
        1.25,
        0.24,
        size=9,
        color=BLUE,
        bold=True,
    )
    add_text(
        slide,
        "data access  ×  qualified truth  ×  takt-time gap  ×  economic value  ×  design-partner commitment",
        2.0,
        6.52,
        10.2,
        0.3,
        size=11.3,
        color=INK,
        bold=True,
    )
    add_footer(slide, 8)
    return slide


def slide_moat(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    add_brand(slide)
    add_kicker(slide, "Compounding advantage")
    add_title(
        slide,
        "Every qualified inspection task improves the intelligence layer.",
        "The defensible asset is the coupled system of physics, defect truth, uncertainty, and deployment learning.",
    )
    center_x, center_y = 6.665, 4.45
    add_circle(slide, center_x - 0.84, center_y - 0.84, 1.68, NAVY)
    add_text(
        slide,
        "DECISION\nSYSTEM",
        center_x - 0.67,
        center_y - 0.34,
        1.34,
        0.72,
        size=13,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    nodes = [
        (1.0, 2.7, "Physics + geometry", "Hard constraints and acquisition priors", BLUE),
        (9.52, 2.7, "Defect truth", "Qualified labels tied to process failure", TEAL),
        (1.0, 5.05, "Uncertainty", "Calibrated fallback and drift signals", AMBER),
        (9.52, 5.05, "Deployment learning", "Takt time, operators, MES/QMS, economics", BLUE),
    ]
    for x, y, title, body, accent in nodes:
        add_card(slide, x, y, 2.82, 1.27, fill=PAPER, line=LINE)
        add_rect(slide, x, y, 0.08, 1.27, accent)
        add_text(slide, title, x + 0.25, y + 0.2, 2.28, 0.28, size=14, color=INK, bold=True)
        add_text(slide, body, x + 0.25, y + 0.61, 2.28, 0.43, size=9.5, color=MUTED)
        sx = x + 2.82 if x < center_x else x
        sy = y + 0.63
        ex = center_x - 0.9 if x < center_x else center_x + 0.9
        ey = center_y
        add_line(slide, sx, sy, ex, ey, accent, 1.5)
    add_rect(slide, 3.76, 6.34, 5.81, 0.49, PALE_TEAL, radius=True)
    add_text(
        slide,
        "The loop compounds only when each task passes joint acceptance.",
        3.76,
        6.47,
        5.81,
        0.21,
        size=10,
        color=TEAL_DARK,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, 9)
    return slide


def slide_business(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, PAPER)
    add_brand(slide)
    add_kicker(slide, "Business model + route resilience")
    add_title(
        slide,
        "Sell the inspection outcome through two delivery routes.",
        "The integrated route maximizes performance; the software route broadens compatibility and reduces hardware dependence.",
    )
    add_card(slide, 0.68, 2.68, 7.24, 3.55, fill=NAVY, line=NAVY)
    add_pill(
        slide,
        "CHAMPION ROUTE",
        1.02,
        3.02,
        1.64,
        fill=RGBColor(16, 67, 86),
        color=MINT,
    )
    add_text(
        slide,
        "Integrated inline CT system",
        1.02,
        3.62,
        5.78,
        0.54,
        size=27,
        color=WHITE,
        bold=True,
        font=FONT_DISPLAY,
    )
    add_bullet_list(
        slide,
        [
            "System sale or lease",
            "Recurring software, model, and service layer",
            "Co-designed for the highest-value inspection task",
        ],
        1.02,
        4.44,
        5.7,
        1.25,
        color=RGBColor(187, 206, 219),
        size=12.2,
    )
    add_card(slide, 8.25, 2.68, 4.4, 3.55, fill=WHITE, line=LINE)
    add_pill(slide, "INDEPENDENT ALTERNATIVE", 8.58, 3.02, 2.23, fill=PALE_BLUE, color=BLUE)
    add_text(
        slide,
        "AI layer for compatible CT",
        8.58,
        3.61,
        3.55,
        0.63,
        size=24,
        color=INK,
        bold=True,
        font=FONT_DISPLAY,
    )
    add_bullet_list(
        slide,
        [
            "Reconstruction + defect decision software",
            "Per-line or annual subscription",
            "Integration and validation services",
        ],
        8.58,
        4.46,
        3.45,
        1.25,
        size=11.2,
    )
    add_text(
        slide,
        "Shared core: task-specific models · uncertainty · workflow · evidence contracts",
        2.01,
        6.58,
        9.28,
        0.25,
        size=10.3,
        color=TEAL_DARK,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, 10)
    return slide


def slide_google(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)
    add_brand(slide, dark=True)
    add_kicker(slide, "Why Gemini Startup Forum", dark=True)
    add_title(
        slide,
        "We want to turn one inspection model into a learning factory system.",
        "Google can help us design the scalable, multimodal, and governed layer around real-time line-side inference.",
        dark=True,
    )
    blocks = [
        (
            "GEMINI",
            "Multimodal quality copilot",
            "Reason across CT volumes, projections, defect evidence, process logs, and engineering documents.",
            MINT,
        ),
        (
            "VERTEX AI",
            "Reproducible 3D model lifecycle",
            "Evaluate scalable training, experiment tracking, model evaluation, and governance.",
            RGBColor(118, 158, 255),
        ),
        (
            "HYBRID CLOUD",
            "Fleet learning with data boundaries",
            "Keep time-critical inference near the line while aggregating permitted evidence for improvement.",
            AMBER,
        ),
    ]
    for i, (label, headline, body, accent) in enumerate(blocks):
        x = 0.68 + i * 4.18
        add_card(
            slide,
            x,
            2.79,
            3.78,
            2.76,
            fill=RGBColor(14, 48, 72),
            line=RGBColor(44, 82, 105),
        )
        add_text(slide, label, x + 0.27, 3.08, 2.7, 0.24, size=9, color=accent, bold=True)
        add_text(
            slide,
            headline,
            x + 0.27,
            3.55,
            3.05,
            0.67,
            size=20,
            color=WHITE,
            bold=True,
            font=FONT_DISPLAY,
        )
        add_text(slide, body, x + 0.27, 4.48, 3.1, 0.7, size=11.2, color=RGBColor(169, 192, 206))
    add_rect(slide, 1.18, 6.01, 11.0, 0.64, RGBColor(16, 67, 86), radius=True)
    add_text(
        slide,
        "Forum outcome we want: a technically credible Google architecture + mentors who can pressure-test industrial AI deployment.",
        1.38,
        6.18,
        10.6,
        0.29,
        size=11,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(
        slide,
        11,
        dark=True,
        note="Roadmap intent only: ScanLine is evaluating these Google technologies and does not claim a current production deployment on Google Cloud.",
    )
    return slide


def slide_team(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    add_brand(slide)
    add_kicker(slide, "China-based two-person founding team")
    add_title(
        slide,
        "CT systems experience meets medical imaging AI.",
        "A compact technical team spanning inline CT architecture, reconstruction, computer vision, and image analysis.",
    )
    people = [
        (
            "ZL",
            "Zhihua Liang",
            "Co-founder · CT Systems & AI",
            "Previously led AI back-end, reconstruction, and system integration for the DeltaRAY X100 inline CT platform.",
            "Inline CT · reconstruction · system integration",
            BLUE,
        ),
        (
            "JL",
            "Dr. Juan Liu",
            "Co-founder · AI Medical Imaging",
            "PhD researcher contributing image analysis, model development, experimental design, and medical imaging perspective.",
            "Medical imaging · computer vision · experiments",
            TEAL,
        ),
    ]
    for i, (initials, name, role, bio, tags, accent) in enumerate(people):
        x = 0.72 + i * 6.12
        add_card(slide, x, 2.68, 5.72, 3.42, fill=PAPER, line=LINE)
        add_circle(slide, x + 0.35, 3.05, 0.84, accent)
        add_text(
            slide,
            initials,
            x + 0.35,
            3.28,
            0.84,
            0.24,
            size=14,
            color=WHITE,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_text(slide, name, x + 1.48, 3.02, 3.5, 0.4, size=23, color=INK, bold=True)
        add_text(slide, role, x + 1.48, 3.5, 3.65, 0.27, size=10, color=accent, bold=True)
        add_text(slide, bio, x + 0.35, 4.15, 5.02, 0.88, size=13.3, color=INK_SOFT)
        add_rect(slide, x + 0.35, 5.39, 5.02, 0.4, WHITE, radius=True, line=LINE)
        add_text(
            slide,
            tags.upper(),
            x + 0.49,
            5.5,
            4.73,
            0.18,
            size=8,
            color=MUTED,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
    add_text(
        slide,
        "Current operating model: two co-founders first, specialist partners added around the selected pilot.",
        2.21,
        6.55,
        8.9,
        0.25,
        size=10.5,
        color=MUTED,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, 12)
    return slide


def slide_close(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)
    add_brand(slide, dark=True)
    add_kicker(slide, "The next construction", dark=True)
    add_text(
        slide,
        "Turn the open bridge\ninto measured evidence.",
        0.65,
        1.4,
        7.25,
        1.8,
        size=42,
        color=WHITE,
        bold=True,
        font=FONT_DISPLAY,
        line_spacing=0.92,
    )
    add_text(
        slide,
        "We are looking for:",
        0.69,
        3.5,
        2.2,
        0.35,
        size=12,
        color=MINT,
        bold=True,
    )
    asks = [
        ("DESIGN PARTNER", "one part · one defect · one line"),
        ("GOOGLE MENTORS", "industrial AI architecture · validation · go-to-market"),
        ("COMPUTE + PLATFORM", "3D model development · multimodal knowledge · hybrid operations"),
    ]
    for i, (label, body) in enumerate(asks):
        y = 4.02 + i * 0.84
        add_circle(slide, 0.72, y + 0.05, 0.24, TEAL)
        add_text(slide, label, 1.18, y, 1.73, 0.25, size=8.5, color=MINT, bold=True)
        add_text(slide, body, 2.88, y - 0.01, 4.8, 0.31, size=12.5, color=WHITE)
    add_card(
        slide,
        8.42,
        1.18,
        4.21,
        5.21,
        fill=RGBColor(14, 48, 72),
        line=RGBColor(44, 82, 105),
    )
    add_text(
        slide,
        "90-DAY SUCCESS TEST",
        8.81,
        1.67,
        3.45,
        0.27,
        size=9,
        color=MINT,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        "Faster cycle time",
        8.81,
        2.36,
        3.45,
        0.48,
        size=25,
        color=WHITE,
        bold=True,
        font=FONT_DISPLAY,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        "at the same qualified\ndefect threshold",
        8.81,
        3.05,
        3.45,
        0.78,
        size=18,
        color=RGBColor(190, 211, 224),
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_line(slide, 9.12, 4.3, 11.94, 4.3, RGBColor(52, 88, 111), 1)
    add_text(
        slide,
        "scanline.cn",
        8.81,
        4.73,
        3.45,
        0.33,
        size=16,
        color=MINT,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        "be.linkedin.com/in/liangzhihua",
        8.81,
        5.3,
        3.45,
        0.25,
        size=9.5,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        "See every part. Decide before it leaves the line.",
        1.72,
        6.75,
        5.77,
        0.28,
        size=11,
        color=RGBColor(163, 188, 203),
        bold=True,
    )
    add_footer(slide, 13, dark=True)
    return slide


def build():
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    prs.core_properties.title = "ScanLine — Gemini Startup Forum 2026"
    prs.core_properties.subject = "AI-native inline CT inspection"
    prs.core_properties.author = "ScanLine"
    prs.core_properties.keywords = "industrial CT, physics-informed AI, manufacturing"

    slide_cover(prs)
    slide_problem(prs)
    slide_solution(prs)
    slide_ai(prs)
    slide_proof(prs)
    slide_state(prs)
    slide_probe(prs)
    slide_wedge(prs)
    slide_moat(prs)
    slide_business(prs)
    slide_google(prs)
    slide_team(prs)
    slide_close(prs)

    prs.save(OUT)
    print(f"Wrote {OUT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
